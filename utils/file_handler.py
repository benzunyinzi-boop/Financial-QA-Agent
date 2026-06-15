import os, hashlib, base64
from collections import OrderedDict
from functools import lru_cache
from utils.logger_handler import logger

from langchain_core.documents import Document
from langchain_community.document_loaders import PyMuPDFLoader, TextLoader

def get_file_md5_hex(filepath:str):         # 获取文件的md5的十六进制字符串
    if not os.path.exists(filepath):
        logger.error(f"[MD5计算]文件{filepath}不存在")
        return
    if not os.path.isfile(filepath):
        logger.error(f"[MD5计算]路径{filepath}不是文件")
        return

    md5_obj = hashlib.md5()
    chunk_size=4096     # 4KB分片，避免文件过大爆内存
    try:
            with open(filepath, "rb")as f:  # 必须二进制读取
                while chunk:=f.read(chunk_size):
                    md5_obj.update(chunk)

            md5_hex=md5_obj.hexdigest()
            return md5_hex
    except Exception as  e:
        logger.error(f"计算文件{filepath}md5失败，{str(e)}")
        return None

def listdir_with_allowed_type(path:str,allowed_types:tuple[str]):   # 返回文件夹内的文件列表（允许的文件后缀）
    files= []
    if not os.path.isdir(path):
        logger.error(f"[lisdir_with_allowed_type]{path}不是文件夹")
        return allowed_types

    for f in os.listdir(path):
        if f.endswith(allowed_types):
            files.append(os.path.join(path,f))
    return tuple(files)

def pdf_loader(filepath: str, passwd: str = None) -> list[Document]:
    # PyMuPDFLoader 相比 PyPDFLoader 能更好地保留表格、多栏布局和页眉页脚结构
    return PyMuPDFLoader(filepath, password=passwd or "").load()


def pdf_markdown_loader(filepath: str) -> list[Document]:
    """
    PDF → Markdown 加载器（保留章节标题结构，便于 MarkdownHeaderTextSplitter 切分）
    适用：标题层级清晰的 PDF（如规范化的保险条款）
    备注：扫描件 / 排版混乱的 PDF 用纯文本 pdf_loader 反而更好
    """
    try:
        import pymupdf4llm
    except ImportError:
        logger.warning("[PDF→Markdown]缺少 pymupdf4llm，降级为 pdf_loader")
        return pdf_loader(filepath)

    try:
        md_text = pymupdf4llm.to_markdown(filepath)
    except Exception as e:
        logger.warning(f"[PDF→Markdown]{filepath} 转 Markdown 失败：{e}，降级为 pdf_loader")
        return pdf_loader(filepath)

    if not md_text or not md_text.strip():
        return pdf_loader(filepath)

    return [Document(
        page_content=md_text,
        metadata={"source": filepath, "original_format": "markdown"},
    )]


def pptx_loader(filepath: str) -> list[Document]:
    """
    PPT 加载器：每页 slide 一个 Document，提取标题/正文/备注/表格
    PPT 一页就是一个完整语义单元，不再做字符级切分
    """
    try:
        from pptx import Presentation
    except ImportError:
        logger.error("[PPT加载]缺少依赖 python-pptx，请安装：pip install python-pptx")
        return []

    docs: list[Document] = []
    try:
        prs = Presentation(filepath)
    except Exception as e:
        logger.error(f"[PPT加载]打开 {filepath} 失败：{e}")
        return []

    for idx, slide in enumerate(prs.slides):
        slide_no = idx + 1

        # 1) 主体文字（标题、正文、形状内文字）
        text_parts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                tf_text = "\n".join(
                    p.text for p in shape.text_frame.paragraphs if p.text.strip()
                )
                if tf_text.strip():
                    text_parts.append(tf_text)

        # 2) 表格（结构化提取为简单文本）
        table_parts = []
        for shape in slide.shapes:
            if shape.has_table:
                rows = []
                for row in shape.table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    rows.append(" | ".join(cells))
                if rows:
                    table_parts.append("\n".join(rows))

        # 3) 备注（销售话术 PPT 的备注通常是干货）
        notes = ""
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()

        # 拼接
        content_blocks = []
        if text_parts:
            content_blocks.append("\n".join(text_parts))
        if table_parts:
            content_blocks.append("【表格】\n" + "\n\n".join(table_parts))
        if notes:
            content_blocks.append("【备注】\n" + notes)

        page_content = "\n\n".join(content_blocks).strip()
        if not page_content:
            continue

        docs.append(Document(
            page_content=page_content,
            metadata={
                "page": slide_no,           # 与 PyMuPDFLoader 的 page 字段对齐
                "slide_index": slide_no,    # PPT 专属字段
            },
        ))

    logger.info(f"[PPT加载]{filepath} 提取 {len(docs)}/{len(prs.slides)} 页有效内容")
    return docs


def txt_loader(filepath:str)->list[Document]:

    return TextLoader(filepath,encoding="utf-8").load()


# ---------- 图片加载（VLM + 缓存）----------

# 默认提取提示词：让 VLM 同时做 OCR + 结构描述 + 关键信息提取
DEFAULT_IMAGE_PROMPT = (
    "请详细提取这张图片中的所有信息：\n"
    "1）逐字提取所有可见文字（OCR），保留原有格式和换行；\n"
    "2）如果有表格，按行列结构输出，单元格用 | 分隔；\n"
    "3）如果是保险相关单据/条款，重点提取保单号、保额、保费、被保险人、生效日期、条款编号等关键字段；\n"
    "4）最后用一两句话总结图片的核心内容。\n"
    "只输出提取结果，不要添加额外说明。"
)


@lru_cache(maxsize=512)
def _vlm_describe_cached(image_md5: str, prompt: str, image_path: str) -> str:
    """
    VLM 描述缓存层（按图片 MD5 缓存）。
    image_path 不参与 hash，仅用于实际调用；不同路径的相同图片复用结果。
    """
    try:
        import dashscope
    except ImportError:
        logger.error("[图片加载]缺少依赖 dashscope")
        return ""

    try:
        with open(image_path, "rb") as f:
            image_b64 = base64.b64encode(f.read()).decode("utf-8")

        # 通过 file:// 协议传本地图片（DashScope SDK 支持）
        ext = os.path.splitext(image_path)[1].lstrip(".").lower() or "png"
        data_url = f"data:image/{ext};base64,{image_b64}"

        resp = dashscope.MultiModalConversation.call(
            model="qwen-vl-max",
            messages=[{
                "role": "user",
                "content": [
                    {"image": data_url},
                    {"text": prompt},
                ]
            }],
            timeout=30,
        )

        if resp.status_code != 200:
            logger.error(f"[图片加载]VLM 调用失败 code={resp.code} msg={resp.message}")
            return ""

        # 兼容不同返回结构
        content = resp.output.choices[0].message.content
        if isinstance(content, list):
            text = "\n".join(c.get("text", "") for c in content if isinstance(c, dict))
        else:
            text = str(content)

        return text.strip()

    except Exception as e:
        logger.error(f"[图片加载]VLM 调用异常 {image_path}：{e}", exc_info=True)
        return ""


def image_loader(filepath: str, prompt: str = None) -> list[Document]:
    """
    图片加载器：用 Qwen-VL-Max 提取文字、表格、关键信息。
    每张图 → 1 个 Document（modality=image_description）。
    内置 MD5 缓存，避免重复调用 VLM。
    """
    if not os.path.exists(filepath):
        logger.error(f"[图片加载]{filepath} 不存在")
        return []

    image_md5 = get_file_md5_hex(filepath)
    if not image_md5:
        return []

    description = _vlm_describe_cached(image_md5, prompt or DEFAULT_IMAGE_PROMPT, filepath)
    if not description:
        logger.warning(f"[图片加载]{filepath} VLM 提取为空，跳过")
        return []

    return [Document(
        page_content=description,
        metadata={
            "source": filepath,
            "modality": "image_description",   # 由 _enrich_metadata 时不会覆盖此值
            "original_image_path": filepath,
            "image_md5": image_md5,
        },
    )]


# ---------- 音频加载（DashScope Paraformer ASR + 缓存）----------

# 简单 LRU 缓存：audio_md5 → tuple[dict] segments
_AUDIO_CACHE: "OrderedDict[str, tuple]" = OrderedDict()
_AUDIO_CACHE_MAX_SIZE = 64

# DashScope Paraformer 支持的音频格式
_SUPPORTED_AUDIO_EXT = ("wav", "mp3", "m4a", "flac", "aac", "opus", "pcm")


def _audio_cache_get(audio_md5: str):
    if audio_md5 in _AUDIO_CACHE:
        _AUDIO_CACHE.move_to_end(audio_md5)
        return _AUDIO_CACHE[audio_md5]
    return None


def _audio_cache_set(audio_md5: str, value: tuple):
    _AUDIO_CACHE[audio_md5] = value
    _AUDIO_CACHE.move_to_end(audio_md5)
    while len(_AUDIO_CACHE) > _AUDIO_CACHE_MAX_SIZE:
        _AUDIO_CACHE.popitem(last=False)


def _transcribe_audio(filepath: str) -> tuple:
    """
    调用 DashScope Paraformer-realtime-v2 转写音频
    返回 tuple[dict]，每个 dict 含 text / begin_time(ms) / end_time(ms)
    异常时返回空 tuple（上层做降级）
    """
    try:
        from dashscope.audio.asr import Recognition, RecognitionCallback
    except ImportError:
        logger.error("[音频加载]缺少 dashscope.audio.asr 模块，跳过 ASR")
        return tuple()

    ext = os.path.splitext(filepath)[1].lstrip(".").lower()
    if ext not in _SUPPORTED_AUDIO_EXT:
        logger.warning(f"[音频加载]不支持的格式：{ext}")
        return tuple()

    # 累积识别结果的回调
    class _Callback(RecognitionCallback):
        def __init__(self):
            self.sentences = []

        def on_open(self):
            pass

        def on_close(self):
            pass

        def on_event(self, result):
            try:
                sentence = result.get_sentence()
                if sentence and sentence.get("text"):
                    self.sentences.append({
                        "text": sentence["text"],
                        "begin_time": sentence.get("begin_time", 0),
                        "end_time": sentence.get("end_time", 0),
                    })
            except Exception as e:
                logger.debug(f"[音频加载]on_event 解析失败：{e}")

    callback = _Callback()
    recognition = Recognition(
        model="paraformer-realtime-v2",
        format=ext,
        sample_rate=16000,    # paraformer-realtime-v2 默认 16kHz
        callback=callback,
    )

    try:
        recognition.start()
        with open(filepath, "rb") as f:
            while True:
                chunk = f.read(3200)    # 100ms@16kHz/16bit/mono = 3200 字节
                if not chunk:
                    break
                recognition.send_audio_frame(chunk)
        recognition.stop()
    except Exception as e:
        logger.error(f"[音频加载]ASR 调用异常 {filepath}：{e}", exc_info=True)
        return tuple()

    return tuple(callback.sentences)


def audio_loader(filepath: str) -> list[Document]:
    """
    音频加载器：用 DashScope Paraformer-realtime-v2 转写
    每个识别出的句子 → 1 个 Document，metadata 含时间戳（用于跳转回放）
    支持 wav / mp3 / m4a / flac / aac / opus / pcm
    内置 LRU 缓存（按音频 MD5），避免重复昂贵 ASR 调用
    """
    if not os.path.exists(filepath):
        logger.error(f"[音频加载]{filepath} 不存在")
        return []

    audio_md5 = get_file_md5_hex(filepath)
    if not audio_md5:
        return []

    # 缓存命中？
    segments = _audio_cache_get(audio_md5)
    if segments is None:
        segments = _transcribe_audio(filepath)
        if segments:
            _audio_cache_set(audio_md5, segments)

    if not segments:
        logger.warning(f"[音频加载]{filepath} ASR 转写为空")
        return []

    docs = []
    for i, seg in enumerate(segments):
        docs.append(Document(
            page_content=seg["text"],
            metadata={
                "source": filepath,
                "modality": "audio_transcript",       # _enrich_metadata 不会覆盖
                "segment_index": i,
                "timestamp_start": seg.get("begin_time", 0) / 1000.0,   # ms → s
                "timestamp_end": seg.get("end_time", 0) / 1000.0,
                "audio_md5": audio_md5,
                "original_audio_path": filepath,
            },
        ))

    logger.info(
        f"[音频加载]{filepath} 转写完成，共 {len(docs)} 句，"
        f"总时长 {docs[-1].metadata['timestamp_end']:.1f}s"
    )
    return docs
